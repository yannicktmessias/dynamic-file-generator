from tkinter import Tk, Frame, Button, Label, Text, Checkbutton, IntVar, filedialog
from pptx import Presentation
from pptxtopdf import convert
# --hidden-import 'comtypes.stream' for auto-py-to-exe

class Application:

    # Define the application's main window content
    def __init__(self, master=None):

        # Define the main window's title
        master.title('Gerador dinâmico de arquivos pptx')
        self.title = Label(master, text='Gerador dinâmico de arquivos pptx')
        self.title["font"] = ("Arial", "14", "bold")
        self.title.pack(pady=10, padx=10)

        # Create container
        self.container = Frame(master)
        self.container.pack(pady=10, padx=10) # show widget, with vertical margin of 10 and horizontal margin of 10

        # Create a button to open the model file
        self.model_open_button = Button(self.container, text='Selecione o arquivo do modelo', command=self.open_model_file)
        self.model_open_button.pack(pady=10)

        # Create a label for the model file path
        self.model_file_path_label = Label(self.container, text='')
        self.model_file_path_label.pack(pady=0)

        # Create a button to open the data file
        self.data_open_button = Button(self.container, text='Selecione o arquivo dos dados', command=self.open_data_file)
        self.data_open_button.pack(pady=10)

        # Create a label for the data file path
        self.data_file_path_label = Label(self.container, text='')
        self.data_file_path_label.pack(pady=0)

        # Create a button to open the destination folder
        self.destination_folder_open_button = Button(self.container, text='Selecione a pasta de destino', command=self.open_destination_folder)
        self.destination_folder_open_button.pack(pady=10)

        # Create a label for the destination folder path
        self.destination_folder_path_label = Label(self.container, text='')
        self.destination_folder_path_label.pack(pady=0)
        
        # Create check button for pdf files
        self.generate_pdf = IntVar()
        self.generate_pdf_check_button = Checkbutton(
            self.container, 
            text = "Gerar arquivos PDF", 
            variable = self.generate_pdf, 
            onvalue = 1, 
            offvalue = 0,
        )
        self.generate_pdf_check_button.pack(pady=10)

        # Create a button to generate the custom files
        self.generate_files_button = Button(self.container, text='Gerar arquivos', command=self.generate_files, state='disabled')
        self.generate_files_button.pack(pady=10)

        # Create a label for the feedback message
        self.feedback_label = Label(self.container, text='')
        self.feedback_label.pack(pady=0)

    # 
    def open_model_file(self):
        file_path = filedialog.askopenfilename(
            title='Selecione o arquivo com o modelo pptx', filetypes=[('Arquivos do PowerPoint', '*.pptx'),]
        )
        if file_path:
            self.feedback_label['text'] = ''
            self.model_file_path_label['text'] = file_path
            if self.data_file_path_label['text'] != '' and self.destination_folder_path_label['text'] != '':
                self.generate_files_button.configure(state='normal')

    #
    def open_data_file(self):
        file_path = filedialog.askopenfilename(
            title='Selecione o arquivo com as variáveis e registros em CSV', filetypes=[('Arquivos CSV', '*.csv'), ('Arquivos de texto', '*.txt')]
        )
        if file_path:
            self.feedback_label['text'] = ''
            self.data_file_path_label['text'] = file_path
            if self.model_file_path_label['text'] != '' and self.destination_folder_path_label['text'] != '':
                self.generate_files_button.configure(state='normal')

    #
    def open_destination_folder(self):
        folder_path = filedialog.askdirectory(
            title='Selecione o diretório de destino'
        )
        if folder_path:
            self.feedback_label['text'] = ''
            self.destination_folder_path_label['text'] = folder_path
            if self.model_file_path_label['text'] != '' and self.data_file_path_label['text'] != '':
                self.generate_files_button.configure(state='normal')

    #
    def generate_files(self):
        data = {}
    
        data_file_path = self.data_file_path_label['text']
        with open(data_file_path, 'r') as file:
            lines = file.read().strip().splitlines()
            
            # Get keys from first line
            keys = lines[0].replace(';', ',').split(',')
            
            # Create an empty list for each key
            for key in keys:
                key = key.strip()
                data[key] = []
            
            # Fill the key lists with values from each line after the first
            for line in lines[1:]:
                values = line.replace(';', ',').split(',')
                for i in range(0, len(keys)):
                    key = keys[i].strip()
                    value = values[i].strip()
                    data[key].append(value)
        
        first_key = next(iter(data))
        number_of_entries = len(data[first_key])        
        model_file_path = self.model_file_path_label['text']
        
        # Update feedback message
        self.feedback_label['text'] = 'Gerando os arquivos pptx...'
        self.feedback_label.update()
        for i in range(number_of_entries):
        
            # Open pptx model file
            prs = Presentation(model_file_path)
                        
            # Replace key for value for each key
            for key in data.keys():
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text == f'<{key}>':
                                    run.text = data[key][i]
            
            generated_file_path = self.destination_folder_path_label['text'] + '/' + data[first_key][i] + '.pptx'
            prs.save(generated_file_path)
            
        
        if self.generate_pdf.get():
            # Update feedback message
            self.feedback_label['text'] = 'Gerando os arquivos PDF...'
            self.feedback_label.update()
            
            # Generate PDF files
            convert(self.destination_folder_path_label['text'], self.destination_folder_path_label['text'])
        
        # Update feedback message
        self.feedback_label['text'] = 'Sucesso!'
        self.feedback_label.update()

# Create the main window (root)
window = Tk()
window.config(padx=10, pady=10)
# window.geometry("500x650")

# Use application on root
Application(window)

# Run the Tkinter event loop
window.mainloop()
